from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from models import GroupedResult

# 색상 정의
COLORS = {
    "title_bg": RGBColor(0x1A, 0x1A, 0x3E),
    "title_text": RGBColor(0xFF, 0xFF, 0xFF),
    "stage_fill": RGBColor(0xE3, 0xF2, 0xFD),
    "stage_border": RGBColor(0x19, 0x76, 0xD2),
    "stage_text": RGBColor(0x0D, 0x47, 0xA1),
    "arrow": RGBColor(0x90, 0xCA, 0xF9),
    "story_text": RGBColor(0x33, 0x33, 0x33),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "black": RGBColor(0x00, 0x00, 0x00),
    "header_bg": RGBColor(0x1E, 0x3A, 0x5F),
    "header_text": RGBColor(0xFF, 0xFF, 0xFF),
}

ROW_COLORS = {
    "기능": RGBColor(0xE8, 0xF5, 0xE9),
    "사양": RGBColor(0xE3, 0xF2, 0xFD),
    "서비스": RGBColor(0xFF, 0xF3, 0xE0),
    "사업요소": RGBColor(0xF3, 0xE5, 0xF5),
}

LABEL_COLORS = {
    "기능": RGBColor(0x2E, 0x7D, 0x32),
    "사양": RGBColor(0x15, 0x65, 0xC0),
    "서비스": RGBColor(0xE6, 0x51, 0x00),
    "사업요소": RGBColor(0x6A, 0x1B, 0x9A),
}

FONT_NAME = "맑은 고딕"


def _set_font(run, size_pt=10, bold=False, color=None, font_name=FONT_NAME):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = font_name
    if color:
        run.font.color.rgb = color


def _add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def _add_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["arrow"]
    shape.line.fill.background()
    return shape


def generate_ppt(grouped: GroupedResult, output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    stages = grouped.stages
    num_stages = len(stages) if stages else 1

    # --- Title ---
    title_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.7),
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLORS["title_bg"]
    title_shape.line.fill.background()
    tf = title_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "User Story Map - 경험 구현 요소"
    _set_font(run, 20, bold=True, color=COLORS["title_text"])

    # --- Stage Flow ---
    margin_left = Inches(0.5)
    stage_area_width = Inches(12.333)
    stage_width = int(stage_area_width / num_stages * 0.7)
    gap = int(stage_area_width / num_stages * 0.3)
    stage_height = Inches(0.65)
    stage_top = Inches(1.0)
    arrow_width = int(gap * 0.6)
    arrow_height = Inches(0.3)

    stage_positions = []  # (left, center_x) for each stage

    for i, stage in enumerate(stages):
        left = int(margin_left + i * (stage_width + gap))
        stage_positions.append((left, left + stage_width // 2))

        # Stage box
        shape = _add_rounded_rect(
            slide, left, stage_top, stage_width, stage_height,
            COLORS["stage_fill"], COLORS["stage_border"]
        )
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = stage.title
        _set_font(run, 11, bold=True, color=COLORS["stage_text"])

        # Arrow (except after last stage)
        if i < num_stages - 1:
            arrow_left = left + stage_width + int((gap - arrow_width) / 2)
            arrow_top = int(stage_top + (stage_height - arrow_height) / 2)
            _add_arrow(slide, arrow_left, arrow_top, arrow_width, arrow_height)

    # --- User Stories ---
    story_top = Inches(1.85)

    for i, stage in enumerate(stages):
        sid = stage.id
        group = grouped.groups.get(sid)
        if not group or not group.user_stories:
            continue

        left = stage_positions[i][0]
        tb = slide.shapes.add_textbox(left, story_top, stage_width, Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True

        for j, story in enumerate(group.user_stories):
            if j > 0:
                p = tf.add_paragraph()
            else:
                p = tf.paragraphs[0]
            p.space_before = Pt(2)
            p.space_after = Pt(2)
            run = p.add_run()
            run.text = f"• {story}"
            _set_font(run, 8, color=COLORS["story_text"])

    # --- Experience Elements Table ---
    categories = ["기능", "사양", "서비스", "사업요소"]
    table_top = Inches(3.3)
    table_left = Inches(0.3)
    table_width = Inches(12.733)
    col_label_width = Inches(1.2)
    col_stage_width = int((table_width - col_label_width) / num_stages)

    rows = len(categories) + 1  # header + 4 categories
    cols = num_stages + 1  # label + stages

    table_shape = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, Inches(3.0))
    table = table_shape.table

    # Column widths
    table.columns[0].width = col_label_width
    for c in range(1, cols):
        table.columns[c].width = col_stage_width

    # Header row
    header_cell = table.cell(0, 0)
    header_cell.text = ""
    header_cell.fill.solid()
    header_cell.fill.fore_color.rgb = COLORS["header_bg"]

    for i, stage in enumerate(stages):
        cell = table.cell(0, i + 1)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["header_bg"]
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = stage.title
        _set_font(run, 9, bold=True, color=COLORS["header_text"])

    # Category rows
    for r, cat in enumerate(categories):
        row_idx = r + 1

        # Label cell
        label_cell = table.cell(row_idx, 0)
        label_cell.fill.solid()
        label_cell.fill.fore_color.rgb = ROW_COLORS[cat]
        label_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = label_cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = cat
        _set_font(run, 10, bold=True, color=LABEL_COLORS[cat])

        # Stage cells
        for i, stage in enumerate(stages):
            cell = table.cell(row_idx, i + 1)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_COLORS[cat]
            cell.vertical_anchor = MSO_ANCHOR.TOP

            group = grouped.groups.get(stage.id)
            if group:
                items = getattr(group.elements, cat, [])
                tf = cell.text_frame
                tf.word_wrap = True
                for j, item_text in enumerate(items):
                    if j > 0:
                        p = tf.add_paragraph()
                    else:
                        p = tf.paragraphs[0]
                    p.space_before = Pt(1)
                    p.space_after = Pt(1)
                    run = p.add_run()
                    run.text = item_text
                    _set_font(run, 8, color=COLORS["black"])

    # Row heights
    for r in range(rows):
        table.rows[r].height = Inches(0.6) if r == 0 else Inches(0.6)

    prs.save(output_path)
